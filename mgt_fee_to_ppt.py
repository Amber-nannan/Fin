# %% 
from Function_mod.basic_creits_date_sql import *
import matplotlib.pyplot as plt
import seaborn as sns
import plotly.subplots as sp
import plotly.graph_objs as go
import plotly.express as px
import plotly.offline as pyo
import warnings
import webbrowser
from pptx import Presentation
from pptx.util import Cm
from pptx.util import Pt
from pptx.dml.color import RGBColor
from io import BytesIO
import pptx_ea_font
warnings.filterwarnings("ignore", category=UserWarning)
plt.rcParams['font.sans-serif'] = ['SimHei']  # 设置中文字体
plt.rcParams['axes.unicode_minus'] = False 

# %%  获取一些基础信息和需要的值
def get_basic_info():
    # REITs_type_L1、2的字典
    db = 'dbconghua'
    engine = create_engine('mysql+pymysql://conghuadb_user:123456@101.132.162.44:3306/'+db)
    con = engine.connect()
    sql_w = f"select ID,reit_code,REITs_name,REITs_type_L1,REITs_type_L2 from dbconghua.basic_info"
    basic_info = pd.read_sql(f"{sql_w}", con, index_col='ID')
    name_dict = dict(zip(basic_info['reit_code'],basic_info['REITs_name']))
    return basic_info, name_dict

def get_mkt_cap_dict():
    # 市值字典(按照2023Q3 market_cap_float划分)
    mkt_cap_df=get_creits_sql_data(['market_cap_float'],'fin_data',False)
    mkt_cap_df = mkt_cap_df.pivot(index=['rp_date','reit_code'], columns='factor_name', values='data').reset_index()

    temp_df = mkt_cap_df[mkt_cap_df['rp_date'] == '2023-09-30'].copy()
    temp_df['market_cap_type'] = pd.qcut(temp_df['market_cap_float'], q=3, labels=['小市值', '中市值', '大市值'])
    mkt_cap_dict=dict(zip(temp_df['reit_code'],temp_df['market_cap_type']))
    return mkt_cap_dict

def get_df(factor_list,tn,trans_dict,basic_info,mkt_cap_dict):
    df_ = get_creits_sql_data(factor_list,tn,False)
    df_=df_.pivot(index=['rp_period','reit_code'], columns='factor_name', values='data').reset_index()
    df_=df_[['rp_period','reit_code']+factor_list]  # 按factor_list排列列的顺序
    df_.rename(columns=trans_dict, inplace=True)
    df_=pd.merge(df_,basic_info,on='reit_code',how='left')
    df_ = df_.fillna(0)
    df_['sum'] = df_.select_dtypes(include='float64').sum(axis=1)
    df_['市值'] = df_['reit_code'].map(mkt_cap_dict)
    df_.sort_values(by=['REITs_type_L1','REITs_type_L2'],inplace=True)
    if '托管费占总结构成本' in df_.columns:
        df_ = df_[df_['托管费占总结构成本'] != 1]
    return df_

def get_mean_values(df_,groupby_):
    factors = df_.select_dtypes(include='float64').columns.tolist()
    df_grouped = df_.groupby(groupby_)
    if 'REITs_name' in groupby_:
        df_= pd.concat([df_grouped[['REITs_type_L1','REITs_type_L2','rp_period']].first(),
                df_grouped[factors].mean()], axis=1).reset_index().sort_values(by='sum',ascending=False)
    elif 'REITs_type_L2' in groupby_:
        df_= pd.concat([df_grouped[['REITs_type_L1','REITs_type_L2','rp_period']].first(),
                df_grouped[factors].mean()], axis=1).reset_index(drop=True).sort_values(by=['REITs_type_L1','REITs_type_L2'])
    elif 'REITs_type_L1' in groupby_:
        df_= df_grouped[factors].mean().reset_index().sort_values(by='REITs_type_L1')
    elif '市值' in groupby_:
        df_= df_grouped[factors].mean().reset_index().sort_values(by='市值', key=lambda x: x.map({'小市值': 1, '中市值': 2, '大市值': 3}))
    return df_

# %% 画图d的设置
color_dict = {
    '产权类': '#F30C7F',
    '经营权类': '#00CC99',
    '产业园': '#1CADE4',
    '保障性住房': '#F89D58',
    '物流仓储': '#AB63FA',
    '市政生态': '#68E622',
    '收费公路': '#018885',
    '能源': '#636EFA',
    # ['#0DF9FF','#63E6E1','#63C1E3','#6196E2','#5E5CE4','#7542A4']
    # [ "#63E6E1", '#5F8FE4',"#5E5CE4"]
    '固定管理费':'#63E6E1',
    '浮动管理费':'#5F8FE4',
    '托管费':'#5E5CE4',

    '固定管理费/nav':'#63E6E1',
    '浮动管理费/nav':'#5F8FE4',
    '托管费/nav':'#5E5CE4',

    '固定管理费占总收入':'#63E6E1',
    '浮动管理费占总收入':'#5F8FE4',
    '托管费占总收入':'#5E5CE4',
    #-----------------------------------------
    '基金和计划管理人管理费':"#63E6E1",
    '外部机构管理费':'#5F8FE4',
    '外部机构管理费净运营费用':'#7542A4',

    '基金和计划管理人管理费占总收入':'#63E6E1',
    '外部机构管理费占总收入':'#5F8FE4',
    '外部机构管理费净运营费用占总收入':'#7542A4',

    '基金和计划管理人管理费/nav':'#63E6E1',
    '外部机构管理费/nav':'#5F8FE4',
    '外部机构管理费净运营费用/nav':'#7542A4',
    # ------------------------------------------
    '小市值':'#63E6E1',
    '中市值':'#5F8FE4',
    '大市值':'#5E5CE4',

    '固定管理费占总结构成本':"#63E6E1",
    '浮动管理费占总结构成本':'#5F8FE4',
    '托管费占总结构成本':'#5E5CE4',

    '基金和计划管理人管理费占比':'#63E6E1',
    '外部机构管理费占比':'#5F8FE4',
    '外部机构管理费净运营费用占比':'#7542A4',
}

style_dict = {
    'dark':dict(
            template='plotly_dark',
            layout_dict = dict(plot_bgcolor='#262626',
                          paper_bgcolor='#262626',
                          font=dict(size=14,color='white'),
                          ),
            yaxes_dict = dict(linecolor='Grey',gridcolor='#5E5E5E',linewidth = 0.1),
            avgline_color = 'white',
            text_bg = '#EAEAEA',
            text_cl = 'black',       
            ),
    'light':dict(
            template='plotly_white',
            layout_dict = dict(plot_bgcolor='#FFFFFF',
                          paper_bgcolor='#FFFFFF',
                          font=dict(size=14,color='black'),
                          ),
            yaxes_dict = dict(linecolor='Grey',gridcolor='#EAEAEA',linewidth = 0.1),
            avgline_color = '#D94B19',
            text_bg = 'black',
            text_cl = 'white',
            text_label = 'black'
            )
    }

def set_global_bar_format(fig):
    fig.update_layout(height=500,width=600, title=dict(y=0.95),
                      legend=dict(orientation='h', x=0, y=1.15))
    fig.update_traces(texttemplate='%{y:.2%}', textposition='inside',
                      textangle=0, textfont=dict(size=12), hovertemplate='=%{y:.2%}')
    fig.update_yaxes(tickformat=".1%")
    fig.update_yaxes(style_dict[style]['yaxes_dict'])
    fig.update_layout(style_dict[style]['layout_dict'])

def set_global_hist_format(fig):
    fig.update_layout(style_dict[style]['layout_dict'],
                    legend=dict(orientation='h',x=0,y=-0.1),
                    width=500,height=500,title=dict(x=0.5))
    # fig.update_xaxes(tickformat=",")
    fig.update_yaxes(gridcolor='#5E5E5E')

def set_global_line_format(fig):
    fig.update_traces(texttemplate='%{y:,}',hovertemplate='=%{y:,}')
    fig.update_layout(height=700,width=1400, legend=dict(orientation='h'))
    fig.update_yaxes(style_dict[style]['yaxes_dict'])
    fig.update_xaxes(style_dict[style]['yaxes_dict'])
    fig.update_layout(style_dict[style]['layout_dict'])

# %% 基础的画图
def barplot(df_,x,y,rp_period,title):
    """
    x: 可选'REITs_name','REITs_type_L1','REITs_type_L2'
    y: 取值为列表，如['固定管理费/nav','浮动管理费/nav']
    """
    df1=df_[df_['rp_period']==rp_period]
    fig=px.bar(df1,x=x,y=y,barmode='stack',title=title,labels={x: '', 'variable': '', 'value': '', }, color_discrete_map=color_dict)
    set_global_bar_format(fig)
    return fig

def histplot(df_,x,rp_period,fig_list):
    df1=df_[df_['rp_period']==rp_period]
    factor_df = df1.select_dtypes(include='float64').drop(columns='sum')
    factor_df = pd.concat([factor_df,df1[x]],axis=1)

    legend_list= df1[x].unique().tolist()  # 设置图例位置,且按照固定顺序排列
    if '产权类' in legend_list:
        factor_df = factor_df.sort_values(by=x, key = lambda x: x.map({'产权类':1,'经营权类':2}))
    elif '产业园' in legend_list:
        factor_df = factor_df.sort_values(by=x, key = lambda x: x.map({'产业园':1,'保障性住房':2,'物流仓储':3,'市政生态':4,'收费公路':5,'能源':6}))
    elif '小市值' in legend_list:
        factor_df = factor_df.sort_values(by=x, key = lambda x: x.map({'小市值':1,'中市值':2,'大市值':3}))
    
    y = df_.select_dtypes(include='float64').columns.tolist()
    y = [variable for variable in y if variable != 'sum']

    for y_i in y:
        fig = px.histogram(factor_df[y_i], color_discrete_map=color_dict, 
                           labels={'variable': '', 'value': '','color':''}, nbins=10, marginal='box', 
                           title=y_i, barmode='stack', color=factor_df[x])
        set_global_hist_format(fig)
        if y_i == y[2]:
            fig.update_layout(legend=dict(orientation='v',x=1,y=0.5,font=dict(size=16)),width=550) 
        else:
            fig.update_layout(showlegend=False) 
        fig_list.append(fig)
    return fig_list

def multi_bar_and_hist_plot(df_,rp_period,title_key_field,fig_list):
    y = df_.select_dtypes(include='float64').columns.tolist()
    y = [variable for variable in y if variable != 'sum']
    
    x = 'REITs_name'
    title= f'<b>各REITs{title_key_field}（{rp_period}）<b>'
    fig=barplot(df_,x,y,rp_period,title)
    fig.update_layout(height=700,width=1400,legend=dict(orientation='h', x=0.3, y=1.05))
    fig_list.append(fig)

    title_key_field = title_key_field.strip('——按固定和浮动分')
    for x in ['REITs_type_L1','REITs_type_L2','市值']:
        title1 = f'<b>按{x}分类的{title_key_field}（{rp_period}）<b>'
        bar_fig = barplot(get_mean_values(df_,groupby_=[x,'rp_period']),x,y,rp_period,title1)
        fig_list.append(bar_fig)
        fig_list = histplot(df_,x,rp_period,fig_list)
    return fig_list

# %% 关于diff的两个图（中报和年报才有）
def diff_plot(div_type,tn,rp_period):
    y_values = [
        'diff',
        f'diff_to_{div_type}',
        f'total_mgt_fee_to_{div_type}',
    ]
    title_field = '/nav' if div_type=='nav' else '占总收入'
    df_ = get_creits_sql_data(y_values, tn, False)
    df_ = df_.pivot(index=['rp_period', 'reit_code'],columns='factor_name', values='data').reset_index()
    df_ = df_[df_['rp_period']==rp_period]
    df_.dropna(how='any', inplace=True)  # drop了一条nan的数据
    df_['REITs_name'] = df_['reit_code'].map(name_dict)
    df_[f'adj_total_mgt_fee_to_{div_type}'] = df_[f'total_mgt_fee_to_{div_type}']+df_[f'diff_to_{div_type}']
    # 这几个股票有黑账
    temp_df = df_[abs(df_['diff']) > 0.1]
    temp_df = temp_df[temp_df['REITs_name'] != '越秀高速']
    temp_df = temp_df[temp_df['REITs_name'] != '红土盐港']  # 这两个中期有diff，但是年报没有diff

    fig = go.Figure()
    fig.add_trace(go.Scatter(x=df_['diff'], y=df_[f'total_mgt_fee_to_{div_type}'], mode='markers', name=f'账面对外支付总报酬{title_field}',
                            marker=dict(size=10,color='#5F8FE4')))
    fig.add_hline(y=df_[f'total_mgt_fee_to_{div_type}'].mean(),line=dict(color='#5F8FE4', width=2))
    fig.add_trace(go.Scatter(x=temp_df['diff'], y=temp_df[f'adj_total_mgt_fee_to_{div_type}'], mode='markers+text', name=f'实际对外支付总报酬{title_field}',
                                text=temp_df['REITs_name'],
                                textposition='top center',  # 指定文本标签的位置
                                marker=dict(size=10,color='#63E6E1')))  # 指定文本标签的位置

    fig.update_layout(title=f'<b>外部机构管理费净运营费用与总结构成本{title_field}散点图（{rp_period}）<b>',
                        xaxis_title='外部机构管理费净运营费用',
                        yaxis_title=f'总结构成本{title_field}',
                        height=700, width=1200, margin=dict(b=150), # 增加底部外边距（以像素为单位）
                        legend=dict(orientation='h', x=0.32, y=1.1),
                        paper_bgcolor='#262626',
                        plot_bgcolor='#262626',
                        font=dict(size=12,color='white')
                        )

    fig.add_annotation(
        text='''注：账面对外支付总报酬=固定管理费+浮动管理费
        <br>      实际对外支付总报酬=基金管理人管理费+计划管理人管理费+外部机构管理费''',
        xref="paper",
        yref="paper",
        x=0, y=-0.3,
        align="left",
        showarrow=False,
        )
    fig.update_xaxes(dict(linecolor='Grey', gridcolor='#5E5E5E',
                     tickformat=",", zerolinecolor='#5E5E5E',))
    fig.update_yaxes(dict(linecolor='Grey', gridcolor='#5E5E5E',
                     tickformat=".1%", zerolinecolor='#5E5E5E',))
    return fig

def plots_related_to_abs_mgt_fee(trans_dict,tn,rp_period,title_key_field):
    fig_list=[]
    factor_list = list(trans_dict.keys())
    df_ = get_df(factor_list, tn, trans_dict, basic_info, mkt_cap_dict).sort_values(by='sum', ascending=False)
    # 柱状图
    y=list(trans_dict.values())
    title = f'<b>各REITs管理费绝对值——{title_key_field}（{rp_period}）<b>'
    fig=barplot(df_,'REITs_name',y,rp_period,title)
    fig.update_yaxes(tickformat=',')
    fig.update_layout(height=700,width=1400,legend=dict(orientation='h', x=0.3, y=1.05))
    fig_list.append(fig)
    
    # 饼图
    mean_value = list(df_[y].mean())
    pie_value = [value/sum(mean_value) for value in mean_value]
    fig = px.pie(values= pie_value,names=y,width=500,height=500)
    fig.update_layout(style_dict[style]['layout_dict'],legend=dict(orientation='h', y=1.25))
    fig.update_traces(marker=dict(colors=[color_dict[y_i] for y_i in y]))
    fig_list.append(fig)
    
    # 直方图
    for y_i in y:
        fig = px.histogram(df_[y_i],color_discrete_map=color_dict,title=y_i,
                        labels={'variable': '', 'value': ''},nbins=10,marginal='box')
        set_global_hist_format(fig)
        fig.update_layout(showlegend=False)
        fig_list.append(fig)
    return fig_list

def plots_related_to_rev_or_nav(trans_dict_list,tn,rp_period,title_key_field):
    fig_list=[]
    trans_dict=trans_dict_list[1]
    df_ = get_df(list(trans_dict.keys()), tn, trans_dict, basic_info, mkt_cap_dict).sort_values(by='sum', ascending=False)
    title = f'<b>各REITs{title_key_field}——按不同主体分（{rp_period}）<b>'
    fig = barplot(df_,'REITs_name',list(trans_dict.values()),rp_period,title)
    fig.update_layout(height=700,width=1400,legend=dict(orientation='h', x=0.3, y=1.05))
    fig_list.append(fig)
    
    trans_dict=trans_dict_list[0]
    df_ = get_df(list(trans_dict.keys()), tn, trans_dict, basic_info, mkt_cap_dict).sort_values(by='sum', ascending=False)
    fig_list = multi_bar_and_hist_plot(df_, rp_period, title_key_field+'——按固定和浮动分', fig_list)
    fig_list[0],fig_list[1] = fig_list[1],fig_list[0]
    return fig_list

def plots_related_to_prop(trans_dict,tn,rp_period):
    fig_list=[]
    factor_list = list(trans_dict.keys())
    df_ = get_df(factor_list, tn, trans_dict, basic_info, mkt_cap_dict)
    if '占总结构成本' in list(trans_dict.values())[0]:
        df_ = df_.sort_values(by=['固定管理费占总结构成本', '浮动管理费占总结构成本'])
        fig_list = multi_bar_and_hist_plot(df_, rp_period, '管理费占总结构成本', fig_list)
    elif '占比' in list(trans_dict.values())[0]:
        df_ = df_.sort_values(by='sum', ascending=False)
        fig_list = multi_bar_and_hist_plot(df_, rp_period, '不同主体管理费占比', fig_list) 
    return fig_list

# %% 季节
def reits_seasonal_patterns(selected_reits,ttm):
    '''
    ttm : bool
    '''
    trans_dict = {
        'mgt_fee_fm_plus_pm': '基金和计划管理人管理费',
        'mgt_fee_for_ext_org': '外部机构管理费',
        'custody_fee': '托管费'
    }
    factor_list = list(trans_dict.keys())
    factor_list_c = list(trans_dict.values())

    df_ = get_df(factor_list, 'fund_mgt_fee_q', trans_dict, basic_info,mkt_cap_dict)
    temp = df_[df_['REITs_name'] == selected_reits].sort_values(by='rp_period')
    if ttm:
        for factor in factor_list_c:
            temp[f'{factor}']=temp[factor].rolling(window=4).sum() 
        fig_title = f'{selected_reits}管理费季度变化（ttm）'
    else :
        fig_title = f'{selected_reits}管理费季度变化'
    fig = px.bar(temp[['rp_period']+factor_list_c], x='rp_period', y=factor_list_c,
                 labels={'variable': '', 'rp_period': ''}, color_discrete_map=color_dict,
                 title=fig_title, width=1400, height=700)
    set_global_bar_format(fig)
    fig.update_yaxes(tickformat=",")
    fig.update_traces(texttemplate='%{y:,}',hovertemplate='=%{y:,}')
    return fig

def reits_L2_seasonal_patterns(selected_L2,selected_factor):
    '''
    selected_factor: 可选择 '总结构成本'、'基金和计划管理人管理费'、'外部机构管理费'、'托管费'
    '''
    trans_dict = {
        'mgt_fee_fm_plus_pm': '基金和计划管理人管理费',
        'mgt_fee_for_ext_org': '外部机构管理费',
        'custody_fee': '托管费'
    }
    factor_list = list(trans_dict.keys())
    df_ = get_df(factor_list, 'fund_mgt_fee_q', trans_dict, basic_info,mkt_cap_dict)
    df_.rename(columns={'sum':'总结构成本'},inplace=True)
    temp=df_[df_['REITs_type_L2']==selected_L2]
    fig = px.line(temp, x='rp_period', y=selected_factor, color='REITs_name',
                  color_discrete_sequence=px.colors.qualitative.Alphabet,
                title=f'{selected_L2}{selected_factor}折线图', labels={'rp_period': '','REITs_name':''})
    set_global_line_format(fig)
    return fig

# style = 'dark'
# basic_info,name_dict = get_basic_info()
# mkt_cap_dict=get_mkt_cap_dict()
# reits_seasonal_patterns('越秀高速')
# reits_L2_seasonal_patterns('产业园','总结构成本')
# %%  存为ppt
def A_slide_layout_1(prs,fig_list):
    '''
    这种布局为2*1，上下两行可以放两个很宽的bar
    '''
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    image_streams=[]
    for fig in fig_list:
        image_stream = BytesIO()
        fig.write_image(image_stream, format="png", scale=3)
        image_stream.seek(0)
        image_streams.append(image_stream)

    left = 3
    top = [0.4, 9.5]
    width = 19.6
    height = 9.1
    for i, image_stream in enumerate(image_streams):
        slide.shapes.add_picture(image_stream, Cm(left), Cm(top[i]), Cm(width), Cm(height))
    return prs

def A_slide_layout_2(prs,fig_list,slide_title):
    '''
    这种布局为 2*4
    '''
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    image_streams=[]
    for fig in fig_list:
        image_stream = BytesIO()
        fig.write_image(image_stream, format="png", scale=3)
        image_stream.seek(0)
        image_streams.append(image_stream)
    
    left = [0, 7, 13, 19]*2
    top = [2.8]*4 + [10.8]*4
    width = [7, 6.4, 6.4, 6.4]*2
    height = 7

    for i, image_stream in enumerate(image_streams):
        slide.shapes.add_picture(image_stream, Cm(left[i]), Cm(top[i]), Cm(width[i]), Cm(height))
    
        # 加标题
    text_box = slide.shapes.add_textbox(left=Cm(1), top=Cm(1), width=Cm(23), height=Cm(0.85))
    text_frame = text_box.text_frame
    text_frame.text = slide_title
    p = text_frame.paragraphs[0]

    for run in p.runs:
        run.font.bold = True
        pptx_ea_font.set_font(run, '等线')
        run.font.size = Pt(14)
        run.font.color.rgb = RGBColor(255, 255, 255)

    return prs

def A_slide_layout_3(prs,fig_list,slide_title):
    '''
    这种布局为 3*4 
    '''
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    image_streams=[]
    for fig in fig_list:
        image_stream = BytesIO()
        fig.write_image(image_stream, format="png", scale=3)
        image_stream.seek(0)
        image_streams.append(image_stream)
    
    left = [1.5, 8.8, 13.9, 19]*3
    top = [2.5] *4 +[7.8] *4 +[13.2] *4
    width = [7.2, 5.5, 5.5, 5.8]*3
    height = [5.7] * 12

    for i, image_stream in enumerate(image_streams):
        slide.shapes.add_picture(image_stream, Cm(left[i]), Cm(top[i]), Cm(width[i]), Cm(height[i]))
    
    # 加标题
    text_box = slide.shapes.add_textbox(left=Cm(1.5), top=Cm(1), width=Cm(23), height=Cm(0.85))
    text_frame = text_box.text_frame
    text_frame.text = slide_title
    p = text_frame.paragraphs[0]

    for run in p.runs:
        run.font.bold = True
        pptx_ea_font.set_font(run, '等线')
        run.font.size = Pt(14)
        run.font.color.rgb = RGBColor(255, 255, 255)
    return prs

def A_slide_layout_4(prs,fig_list):
    '''
    这种布局为 2*1，但第二行额外被分为三列
    '''
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    image_streams=[]
    for fig in fig_list:
        image_stream = BytesIO()
        fig.write_image(image_stream, format="png", scale=3)
        image_stream.seek(0)
        image_streams.append(image_stream)
    
    left = [0,1,8.5,16.5]
    top = [0] +[11.8]*3 
    width = [25.4] + [8]*3
    height = [12] + [6.4]*3

    for i, image_stream in enumerate(image_streams):
        slide.shapes.add_picture(image_stream, Cm(left[i]), Cm(top[i]), Cm(width[i]), Cm(height[i]))
    return prs

def A_slide_layout_5(prs,fig):
    '''
    这种布局为最简单的 1*1
    '''
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    image_stream = BytesIO()   # 将Plotly图形转换为图像
    fig.write_image(image_stream, format="png",scale=3)
    image_stream.seek(0)
    left = Cm(1.2)
    top = Cm(2.3)
    width = Cm(23.5)
    height = Cm(14)
    slide.shapes.add_picture(image_stream, left, top, width, height)
    return prs

def save_prs(fig_list1,fig_list2,fig_list3,fig_list4,fig_list5,tn,rp_period):
    prs = Presentation(r'mgt_fee\mgt_fee_PPT\default_template.pptx')
    prs = A_slide_layout_1(prs,fig_list1[:1]+fig_list1[5:6])
    prs = A_slide_layout_2(prs,fig_list1[1:5]+fig_list1[6:],'管理费两种分解方式下的饼图和直方图')
    prs = A_slide_layout_1(prs,fig_list2[:2])
    prs = A_slide_layout_3(prs,fig_list2[2:],'管理费/nav——按REITs_type_L1、REITs_type_L2、市值分类')
    prs = A_slide_layout_1(prs,fig_list3[:2])
    prs = A_slide_layout_3(prs,fig_list3[2:],'管理费占总收入——按REITs_type_L1、REITs_type_L2、市值分类')
    prs = A_slide_layout_4(prs,fig_list4[:2]+fig_list4[5:6]+fig_list4[9:10])
    prs = A_slide_layout_4(prs,fig_list5[:2]+fig_list5[5:6]+fig_list5[9:10])
    prs = A_slide_layout_5(prs,diff_plot('nav',tn,rp_period))
    prs = A_slide_layout_5(prs,diff_plot('rev',tn,rp_period))
    prs.save(fr"mgt_fee\mgt_fee_PPT\mgt_fee_{rp_period}.pptx")
    
# %%
def main(tn,rp_period):
    fig_list=[]
    # if 'Q' in rp_period:
    #     trans_dict_list = [{
    #         'mgt_fee_fm_plus_pm_to_rev': '基金和计划管理人管理费占总收入',
    #         'mgt_fee_for_ext_org_to_rev': '外部机构管理费占总收入',
    #         'custody_fee_to_rev': '托管费占总收入'
    #     },
    #         {
    #         'mgt_fee_fm_plus_pm_to_nav': '基金和计划管理人管理费/nav',
    #         'mgt_fee_for_ext_org_to_nav': '外部机构管理费/nav',
    #         'custody_fee_to_nav': '托管费/nav'
    #     },
    #         {
    #         'mgt_fee_fm_plus_pm_prop': '基金和计划管理人管理费占比',
    #         'mgt_fee_for_ext_org_prop': '外部机构管理费占比',
    #     }]
    # else:
    #     pass
    # 第一部分绝对值，两张图
    trans_dict_list=[
        {
            'fixed_mgt_fee': '固定管理费',   
            'float_mgt_fee': '浮动管理费',
            'custody_fee': '托管费'
        },
        {
            'mgt_fee_fm_plus_pm':'基金和计划管理人管理费',
            'mgt_fee_for_ext_org':'外部机构管理费',
            'diff':'外部机构管理费净运营费用',
        }
    ]
    fig_list1 = plots_related_to_abs_mgt_fee(trans_dict_list[0],tn,rp_period,'按固定和浮动分')+\
        plots_related_to_abs_mgt_fee(trans_dict_list[1],tn,rp_period,'按不同主体分')

    # 第二部分，/nav和/rev，两张图
    trans_dict_list = [
        {
            'fixed_mgt_fee_to_nav': '固定管理费/nav',   
            'float_mgt_fee_to_nav': '浮动管理费/nav',
            'custody_fee_to_nav': '托管费/nav'
        },
        {
            'mgt_fee_fm_plus_pm_to_nav':'基金和计划管理人管理费/nav',
            'mgt_fee_for_ext_org_to_nav':'外部机构管理费/nav',
            'diff_to_nav':'外部机构管理费净运营费用/nav',
        },
        {
            'fixed_mgt_fee_to_rev': '固定管理费占总收入',
            'float_mgt_fee_to_rev': '浮动管理费占总收入',
            'custody_fee_to_rev': '托管费占总收入'
        },
        {
            'mgt_fee_fm_plus_pm_to_rev':'基金和计划管理人管理费占总收入',
            'mgt_fee_for_ext_org_to_rev':'外部机构管理费占总收入',
            'diff_to_rev':'外部机构管理费净运营费用占总收入',
        }
    ]
    fig_list2 = plots_related_to_rev_or_nav(trans_dict_list[:2],tn,rp_period,'管理费/nav')
    fig_list3 = plots_related_to_rev_or_nav(trans_dict_list[2:],tn,rp_period,'管理费占总收入')

    # 第三部分，比例
    trans_dict_list = [ 
        {
            'fixed_mgt_fee_to_total_struct_costs': '固定管理费占总结构成本',  # 第四部分 一张图
            'float_mgt_fee_to_total_struct_costs': '浮动管理费占总结构成本',
            'custody_fee_to_total_struct_costs': '托管费占总结构成本'
        },   
        {
            'mgt_fee_fm_plus_pm_prop': '基金和计划管理人管理费占比',  # 第五部分 一张图
            'mgt_fee_for_ext_org_prop': '外部机构管理费占比',
            'diff_to_total_mgt_fee': '外部机构管理费净运营费用占比'
        } 
    ]
    fig_list4 = plots_related_to_prop(trans_dict_list[0],tn,rp_period)
    fig_list5 = plots_related_to_prop(trans_dict_list[1],tn,rp_period)
    save_prs(fig_list1,fig_list2,fig_list3,fig_list4,fig_list5,tn,rp_period)


# %% 
style = 'dark'
basic_info,name_dict = get_basic_info()
mkt_cap_dict=get_mkt_cap_dict()
main('fund_mgt_fee_h','2023H1')

# %% 季报和年报对比
def reits_seasonal_patterns(selected_reits,ttm):
    '''
    ttm : bool
    '''
    trans_dict = {
        'mgt_fee_fm_plus_pm': '基金和计划管理人管理费',
        'mgt_fee_for_ext_org': '外部机构管理费',
        'custody_fee': '托管费'
    }
    factor_list = list(trans_dict.keys())
    factor_list_c = list(trans_dict.values())
    df_ = get_df(factor_list, 'fund_mgt_fee_q', trans_dict, basic_info,mkt_cap_dict)
    # 处理1
    if selected_reits in ['建信中关村','国君东久','国君临港']:
        del trans_dict['mgt_fee_for_ext_org']
        trans_dict['diff'] = '外部机构管理费'
        factor_list = list(trans_dict.keys())
        df_1= get_df(factor_list, 'fund_mgt_fee_a', trans_dict, basic_info,mkt_cap_dict)
    else:
        df_1= get_df(factor_list, 'fund_mgt_fee_a', trans_dict, basic_info,mkt_cap_dict)
    # 处理2    
    if selected_reits in ['华夏北京保障房','中金厦门安居','深圳红土安居']:
        selected_indices = df_.query(f"REITs_name == '{selected_reits}' and rp_period == '2022Q4'").index
        float_columns = df_.select_dtypes(include='float')
        df_.loc[selected_indices, float_columns.columns] *= 4/3
        
    temp = df_[df_['REITs_name'] == selected_reits].sort_values(by='rp_period')
    temp1 = df_1[df_1['REITs_name'] == selected_reits].sort_values(by='rp_period')
    # 选择历史累计还是ttm
    if ttm:
        for factor in factor_list_c:
            temp[f'{factor}']=temp[factor].rolling(window=4,min_periods=1).sum() 
        temp.reset_index(inplace=True)
        float_columns = temp.select_dtypes(include='float')
        for idx, row in temp.iterrows():
            if idx == 0:
                temp.loc[idx, float_columns.columns] *= 4
            elif idx == 1:
                temp.loc[idx, float_columns.columns] *= 2
            elif idx == 2:
                temp.loc[idx, float_columns.columns] *= 4/3
        fig_title = f'{selected_reits}管理费季度变化（ttm）'
    else :
        for factor in factor_list_c:
            temp[f'{factor}']=temp[factor].rolling(window=4,min_periods=1).sum() 
        fig_title = f'{selected_reits}管理费季度变化（历史累计）'
    # 画图
    temp = pd.concat([temp,temp1],axis=0)
    fig = px.bar(temp[['rp_period']+factor_list_c], x='rp_period', y=factor_list_c,
                 labels={'variable': '', 'rp_period': ''}, color_discrete_map=color_dict,
                 title=fig_title, width=1400, height=700)
    set_global_bar_format(fig)
    fig.update_yaxes(tickformat=",")
    fig.update_traces(texttemplate='%{y:,}',hovertemplate='=%{y:,}')
    return fig

fig_list=[]
for name in basic_info['REITs_name']:
    fig_list.append(reits_seasonal_patterns(name,True))
html_file = r'mgt_fee\mgt_fee_PPT\ttm_q_compare_with_a.html'
with open(html_file, 'w') as f:
    for i, fig in enumerate(fig_list):
        f.write(fig.to_html(full_html=False, include_plotlyjs='cdn'))
webbrowser.open(html_file)